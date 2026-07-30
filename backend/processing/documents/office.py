"""PPTX and DOCX parsing.

Neither format is really "visual" underneath: both are zipped XML holding real
text frames, real table cells, and -- for PowerPoint -- real chart data and
speaker notes. Reading that structure recovers most of a deck without rendering
a single pixel, which is why this path needs no LibreOffice.

What structure cannot recover is SmartArt, arrow-and-box diagrams, spatial
grouping, and screenshots. Embedded pictures are therefore handed to the vision
model individually, and diagram-heavy decks are better exported to PDF, where
the whole page is rendered.
"""

from __future__ import annotations

import logging
from typing import Iterator, List, Optional

from backend.utils.vision import VisionImage

from .markdown_tables import iter_non_empty, rows_to_markdown_table
from .types import DocumentSource, PageSource

logger = logging.getLogger(__name__)

# Word has no page boundaries without a rendering engine, so sections are cut at
# top-level headings instead. This bounds a section that contains no headings at
# all, so a heading-free report does not become one enormous page.
DOCX_SECTION_CHAR_BUDGET = 6000

_PICTURE_MEDIA_TYPES = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
}


# --------------------------------------------------------------------------
# PowerPoint
# --------------------------------------------------------------------------


def _chart_to_markdown(chart) -> Optional[str]:
    """A native chart's underlying data, as a table.

    This is the reason decks do not need rendering. A vision model reading a
    rendered bar chart estimates the values; python-pptx hands over the exact
    numbers the chart was built from, because they are stored in the file.
    """
    try:
        categories = [str(category) for category in chart.plots[0].categories]
        header = ["Category"] + [
            series.name or f"Series {index + 1}"
            for index, series in enumerate(chart.series)
        ]
        rows = []
        for row_index, category in enumerate(categories):
            row: List[object] = [category]
            for series in chart.series:
                values = list(series.values)
                row.append(values[row_index] if row_index < len(values) else "")
            rows.append(row)
    except Exception as e:  # noqa: BLE001 - chart shapes vary widely
        logger.debug("Chart data extraction failed: %s", e)
        return None

    table = rows_to_markdown_table(rows, header=header)
    if not table:
        return None
    title = ""
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text.strip()
    except Exception:  # noqa: BLE001
        title = ""
    heading = f"Chart: {title}" if title else "Chart"
    return f"**{heading}**\n\n{table}"


def _picture_image(shape) -> Optional[VisionImage]:
    try:
        image = shape.image
        media_type = _PICTURE_MEDIA_TYPES.get((image.ext or "").lower(), "image/png")
        return VisionImage(data=image.blob, media_type=media_type)
    except Exception as e:  # noqa: BLE001 - linked or corrupt pictures
        logger.debug("Picture extraction failed: %s", e)
        return None


def _iter_shapes(shapes):
    """Flatten group shapes, in reading order.

    Shape order in the file is z-order, not reading order, so a slide read
    naively can emit its footer before its title. Sorting by position restores
    the order a human reads in. Groups are recursed into, since a grouped
    diagram's labels are real text that would otherwise be lost.
    """
    try:
        ordered = sorted(
            shapes,
            key=lambda shape: (
                getattr(shape, "top", None) or 0,
                getattr(shape, "left", None) or 0,
            ),
        )
    except Exception:  # noqa: BLE001 - placeholders may lack geometry
        ordered = list(shapes)

    for shape in ordered:
        if getattr(shape, "shape_type", None) is not None and shape.shape_type == 6:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _slide_notes(slide) -> Optional[str]:
    """Speaker notes.

    Worth calling out: these are invisible on the rendered slide, so no amount
    of visual parsing would ever recover them, and they are frequently the most
    useful text in the whole deck for meeting context.
    """
    try:
        if not slide.has_notes_slide:
            return None
        notes = slide.notes_slide.notes_text_frame.text.strip()
    except Exception as e:  # noqa: BLE001
        logger.debug("Notes extraction failed: %s", e)
        return None
    return notes or None


def _slide_title(slide) -> Optional[str]:
    try:
        if slide.shapes.title is not None:
            return (slide.shapes.title.text or "").strip() or None
    except Exception:  # noqa: BLE001 - layouts without a title placeholder
        pass
    return None


def _shape_content(
    shape,
    *,
    title: Optional[str],
    want_images: bool,
) -> tuple[Optional[str], Optional[VisionImage]]:
    """One shape's contribution: a text fragment, an image, or neither.

    Checked in order of specificity. A table and a chart both also report
    ``has_text_frame``, so testing that first would reduce a chart to its axis
    labels and silently drop the data behind it.
    """
    if getattr(shape, "has_table", False):
        rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
        return rows_to_markdown_table(rows), None
    if getattr(shape, "has_chart", False):
        return _chart_to_markdown(shape.chart), None
    if getattr(shape, "shape_type", None) == 13:  # PICTURE
        return None, (_picture_image(shape) if want_images else None)
    if getattr(shape, "has_text_frame", False):
        text = (shape.text_frame.text or "").strip()
        # The title is emitted as the page heading; repeating it in the body
        # just wastes context.
        if text and text != title:
            return text, None
    return None, None


def open_pptx(path: str, *, want_images: bool) -> DocumentSource:
    from pptx import Presentation

    presentation = Presentation(path)
    slides = list(presentation.slides)

    def pages() -> Iterator[PageSource]:
        for index, slide in enumerate(slides):
            title = _slide_title(slide)
            fragments: List[str] = []
            images: List[VisionImage] = []

            for shape in _iter_shapes(slide.shapes):
                fragment, image = _shape_content(
                    shape, title=title, want_images=want_images
                )
                if fragment:
                    fragments.append(fragment)
                if image is not None:
                    images.append(image)

            notes = _slide_notes(slide)
            if notes:
                fragments.append(f"**Speaker notes:** {notes}")

            yield PageSource(
                page_number=index + 1,
                title=title,
                text="\n\n".join(iter_non_empty(fragments)),
                images=images,
            )

    return DocumentSource(page_count=len(slides), pages=pages)


# --------------------------------------------------------------------------
# Word
# --------------------------------------------------------------------------


def _docx_body_blocks(document):
    """Paragraphs and tables in document order.

    python-docx exposes the two as separate collections, which loses their
    interleaving -- a table would end up detached from the paragraph that
    introduces it. Walking the body XML preserves it.
    """
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = document.element.body
    for child in body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            yield Paragraph(child, document)
        elif tag == "tbl":
            yield Table(child, document)


def _docx_images(document, *, want_images: bool) -> List[VisionImage]:
    """Every embedded image in the document.

    Attached to the first section rather than positioned: Word's relationship
    model does not make it cheap to say which heading an inline shape falls
    under, and an image the model sees under the wrong heading is still far
    better than one it never sees.
    """
    if not want_images:
        return []
    images: List[VisionImage] = []
    try:
        for rel in document.part.rels.values():
            if "image" not in rel.reltype:
                continue
            part = rel.target_part
            content_type = getattr(part, "content_type", "") or "image/png"
            if not content_type.startswith("image/"):
                continue
            images.append(VisionImage(data=part.blob, media_type=content_type))
    except Exception as e:  # noqa: BLE001
        logger.debug("DOCX image extraction failed: %s", e)
    return images


class _DocxSectionBuilder:
    """Accumulates blocks into heading-bounded sections."""

    def __init__(self) -> None:
        self.sections: List[PageSource] = []
        self._title: Optional[str] = None
        self._fragments: List[str] = []

    def flush(self) -> None:
        body = "\n\n".join(iter_non_empty(self._fragments))
        if body or self._title:
            self.sections.append(
                PageSource(
                    page_number=len(self.sections) + 1,
                    title=self._title,
                    text=body,
                )
            )
        self._fragments = []
        self._title = None

    def start_section(self, title: str) -> None:
        self.flush()
        self._title = title

    def add(self, fragment: str) -> None:
        self._fragments.append(fragment)
        # Cut a runaway section that never meets another heading, so a
        # heading-free report does not collapse into one enormous page.
        if sum(len(item) for item in self._fragments) >= DOCX_SECTION_CHAR_BUDGET:
            self.flush()


def open_docx(path: str, *, want_images: bool) -> DocumentSource:
    from docx import Document as DocxDocument
    from docx.table import Table

    document = DocxDocument(path)
    builder = _DocxSectionBuilder()

    for block in _docx_body_blocks(document):
        if isinstance(block, Table):
            rows = [[cell.text for cell in row.cells] for row in block.rows]
            builder.add(rows_to_markdown_table(rows))
            continue

        text = (block.text or "").strip()
        if not text:
            continue
        style = (getattr(block.style, "name", "") or "").lower()
        if style.startswith("heading 1") or style.startswith("heading 2"):
            builder.start_section(text)
        else:
            builder.add(text)

    builder.flush()
    sections = builder.sections or [PageSource(page_number=1)]

    images = _docx_images(document, want_images=want_images)
    if images:
        sections[0].images.extend(images)

    def pages() -> Iterator[PageSource]:
        yield from sections

    return DocumentSource(page_count=len(sections), pages=pages)
