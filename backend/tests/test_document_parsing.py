"""Structural parsing, chunking, and the vision merge rules.

Every format is exercised against a real file built in a fixture rather than a
recorded string, so a library upgrade that changes extraction behaviour fails
here instead of silently degrading uploads in production.
"""

from __future__ import annotations

import pytest

from backend.processing.documents import (
    UnsupportedDocumentError,
    chunk_page_content,
    is_rendered_page_format,
    is_vision_only_format,
    open_document,
)
from backend.processing.documents.markdown_tables import rows_to_markdown_table
from backend.processing.documents.types import PageSource
from backend.processing.documents.vision import (
    _strip_code_fence,
    build_figure_prompt,
    build_page_prompt,
    merge_page_content,
)
from backend.utils.vision import VisionImage

# ---------------------------------------------------------------------------
# Fixtures: one real file per format
# ---------------------------------------------------------------------------


@pytest.fixture
def pptx_path(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Q3 Roadmap"
    slide.placeholders[1].text = "Ship search\nHire two engineers"
    slide.notes_slide.notes_text_frame.text = "Budget sign-off still pending."

    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Revenue"
    chart_data = CategoryChartData()
    chart_data.categories = ["EMEA", "APAC"]
    chart_data.add_series("FY25", (1200.0, 850.0))
    second.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(2),
        Inches(6),
        Inches(4),
        chart_data,
    )

    path = tmp_path / "deck.pptx"
    presentation.save(str(path))
    return str(path)


@pytest.fixture
def docx_path(tmp_path):
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_heading("Introduction", level=1)
    document.add_paragraph("This report covers the third quarter.")
    document.add_heading("Findings", level=1)
    document.add_paragraph("Latency improved.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "p95"
    table.cell(1, 1).text = "210ms"

    path = tmp_path / "report.docx"
    document.save(str(path))
    return str(path)


@pytest.fixture
def xlsx_path(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    summary = workbook.active
    summary.title = "Summary"
    summary.append(["Region", "Revenue"])
    summary.append(["EMEA", 1200])
    detail = workbook.create_sheet("Detail")
    detail.append(["Account", "Owner"])
    detail.append(["Acme", "Priya"])

    path = tmp_path / "book.xlsx"
    workbook.save(str(path))
    return str(path)


@pytest.fixture
def pdf_path(tmp_path):
    fitz = pytest.importorskip("fitz")

    document = fitz.open()
    first = document.new_page()
    first.insert_text((72, 100), "Quarterly Review", fontsize=18)
    first.insert_text((72, 140), "Revenue grew 12% in EMEA.", fontsize=11)
    document.new_page().insert_text((72, 100), "Second page content.", fontsize=11)

    path = tmp_path / "doc.pdf"
    document.save(str(path))
    return str(path)


def _pages(path, *, want_images=False):
    source = open_document(path, want_images=want_images)
    return source.page_count, list(source.pages())


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------


def test_pptx_yields_one_page_per_slide_with_titles(pptx_path):
    count, pages = _pages(pptx_path)
    assert count == 2
    assert [page.page_number for page in pages] == [1, 2]
    assert pages[0].title == "Q3 Roadmap"
    assert pages[1].title == "Revenue"


def test_pptx_captures_speaker_notes(pptx_path):
    """Notes are invisible on the rendered slide, so no visual parse would see
    them, yet they are often the most useful text in a deck."""
    _, pages = _pages(pptx_path)
    assert "Budget sign-off still pending." in pages[0].text
    assert "Speaker notes" in pages[0].text


def test_pptx_extracts_exact_chart_values_not_estimates(pptx_path):
    """The reason decks need no rendering: the numbers are in the file.

    A vision model shown a rendered bar chart infers values from pixel heights.
    python-pptx returns what the chart was built from.
    """
    _, pages = _pages(pptx_path)
    body = pages[1].text
    assert "1200" in body
    assert "850" in body
    assert "EMEA" in body and "APAC" in body


def test_pptx_body_does_not_repeat_the_title(pptx_path):
    _, pages = _pages(pptx_path)
    assert pages[0].text.count("Q3 Roadmap") == 0


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------


def test_docx_splits_into_sections_on_headings(docx_path):
    count, pages = _pages(docx_path)
    assert count == 2
    assert [page.title for page in pages] == ["Introduction", "Findings"]


def test_docx_keeps_a_table_with_its_own_section(docx_path):
    """python-docx exposes paragraphs and tables separately, which loses their
    order; the parser walks the body XML so a table stays where it was."""
    _, pages = _pages(docx_path)
    assert "| Metric | Value |" in pages[1].text
    assert "210ms" in pages[1].text
    assert "Metric" not in pages[0].text


# ---------------------------------------------------------------------------
# Spreadsheets and CSV
# ---------------------------------------------------------------------------


def test_xlsx_yields_one_page_per_sheet_named_after_it(xlsx_path):
    count, pages = _pages(xlsx_path)
    assert count == 2
    assert [page.title for page in pages] == ["Summary", "Detail"]
    assert "| EMEA | 1200 |" in pages[0].text
    assert "| Acme | Priya |" in pages[1].text


def test_csv_becomes_a_single_markdown_table(tmp_path):
    path = tmp_path / "data.csv"
    path.write_text("name,role\nPriya,PM\nSam,Eng\n")

    count, pages = _pages(str(path))
    assert count == 1
    assert "| name | role |" in pages[0].text
    assert "| Priya | PM |" in pages[0].text


def test_csv_skips_entirely_blank_rows(tmp_path):
    path = tmp_path / "gappy.csv"
    path.write_text("a,b\n1,2\n,\n3,4\n")

    _, pages = _pages(str(path))
    assert pages[0].text.count("\n") == 3  # header, delimiter, two data rows


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def test_pdf_extracts_text_per_page(pdf_path):
    count, pages = _pages(pdf_path)
    assert count == 2
    assert "Quarterly Review" in pages[0].text
    assert "Second page content." in pages[1].text


def test_pdf_renders_one_image_per_page_only_when_asked(pdf_path):
    _, without = _pages(pdf_path, want_images=False)
    assert all(not page.images for page in without)

    _, with_images = _pages(pdf_path, want_images=True)
    assert all(len(page.images) == 1 for page in with_images)
    assert all(page.images[0].media_type == "image/png" for page in with_images)
    assert all(page.images[0].data.startswith(b"\x89PNG") for page in with_images)


# ---------------------------------------------------------------------------
# Plain text and images
# ---------------------------------------------------------------------------


def test_text_and_markdown_are_single_pages(tmp_path):
    for name, body in [("a.txt", "Plain agenda."), ("b.md", "# Agenda\n\n- One")]:
        path = tmp_path / name
        path.write_text(body)
        count, pages = _pages(str(path))
        assert count == 1
        assert pages[0].text == body


def test_image_upload_has_no_text_but_always_carries_its_image(tmp_path):
    """`want_images=False` is ignored for images on purpose: without the image
    there is nothing at all, and the orchestrator needs to see that."""
    fitz = pytest.importorskip("fitz")
    path = tmp_path / "shot.png"
    pixmap = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 8, 8))
    pixmap.clear_with(200)
    pixmap.save(str(path))

    _, pages = _pages(str(path), want_images=False)
    assert pages[0].text == ""
    assert len(pages[0].images) == 1


def test_unsupported_extension_is_rejected(tmp_path):
    path = tmp_path / "archive.zip"
    path.write_bytes(b"PK\x03\x04")
    with pytest.raises(UnsupportedDocumentError):
        open_document(str(path), want_images=False)


# ---------------------------------------------------------------------------
# Format classification
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "path,rendered,vision_only",
    [
        ("a.pdf", True, False),
        ("a.png", True, True),
        ("a.jpg", True, True),
        ("a.pptx", False, False),
        ("a.docx", False, False),
        ("a.csv", False, False),
    ],
)
def test_format_classification(path, rendered, vision_only):
    assert is_rendered_page_format(path) is rendered
    assert is_vision_only_format(path) is vision_only


# ---------------------------------------------------------------------------
# Markdown tables
# ---------------------------------------------------------------------------


def test_table_escapes_pipes_and_flattens_newlines():
    """An unescaped pipe ends the cell and a newline ends the row, so either
    would silently corrupt every following column."""
    table = rows_to_markdown_table([["h1", "h2"], ["a|b", "c\nd"]])
    assert "a\\|b" in table
    assert "c d" in table
    assert table.count("\n") == 2


def test_table_pads_ragged_rows_to_a_uniform_width():
    table = rows_to_markdown_table([["a", "b", "c"], ["1"]])
    assert table.splitlines()[-1] == "| 1 |  |  |"


def test_table_of_nothing_renders_nothing():
    assert rows_to_markdown_table([]) == ""


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------


def test_a_page_within_the_window_stays_one_chunk():
    """The whole point of the 8192-token model: a page is one retrieval unit."""
    assert chunk_page_content("A short page.") == ["A short page."]


def test_an_empty_page_indexes_nothing():
    assert chunk_page_content("   \n\n  ") == []


def test_an_oversized_page_splits_on_paragraph_boundaries():
    blocks = [f"para {index} " + "x" * 100 for index in range(20)]
    chunks = chunk_page_content("\n\n".join(blocks), budget=500)
    assert len(chunks) > 1
    assert all(len(chunk) <= 500 for chunk in chunks)
    # No paragraph is cut in half.
    assert all(chunk.startswith("para") for chunk in chunks)


def test_an_oversized_table_splits_on_row_boundaries():
    rows = [f"| cell{index} | value{index} |" for index in range(60)]
    chunks = chunk_page_content("\n".join(rows), budget=300)
    assert len(chunks) > 1
    for chunk in chunks:
        for line in chunk.splitlines():
            assert line.startswith("| cell")


def test_a_single_line_longer_than_the_budget_is_cut_bluntly():
    chunks = chunk_page_content("x" * 1000, budget=300)
    assert len(chunks) == 4
    assert all(len(chunk) <= 300 for chunk in chunks)


# ---------------------------------------------------------------------------
# Vision merge rules and prompts
# ---------------------------------------------------------------------------


def test_rendered_page_vision_output_replaces_the_text_layer():
    """The model was given that text and asked to improve on it, so keeping
    both would duplicate every paragraph."""
    page = PageSource(page_number=1, text="raw text layer")
    assert (
        merge_page_content(page, "richer markdown", is_rendered_page=True)
        == "richer markdown"
    )


def test_figure_description_supplements_slide_content():
    """A slide's text frames, tables, chart data and notes were never shown to
    the model, so a picture description is additive rather than a replacement."""
    page = PageSource(page_number=1, text="slide body")
    merged = merge_page_content(page, "a photo of a whiteboard", is_rendered_page=False)
    assert merged == "slide body\n\na photo of a whiteboard"


def test_merge_falls_back_to_structural_when_vision_returned_nothing():
    page = PageSource(page_number=1, text="slide body")
    assert merge_page_content(page, None, is_rendered_page=True) == "slide body"
    assert merge_page_content(page, "  ", is_rendered_page=True) == "slide body"


def test_page_prompt_includes_the_text_layer_when_there_is_one():
    with_text = build_page_prompt(PageSource(page_number=1, text="hello"))
    assert "<extracted_text>" in with_text
    assert "hello" in with_text

    without_text = build_page_prompt(PageSource(page_number=1, text=""))
    assert "<extracted_text>" not in without_text


def test_figure_prompt_names_the_slide_when_it_has_a_title():
    titled = build_figure_prompt(PageSource(page_number=4, title="Revenue"))
    assert 'the slide titled "Revenue"' in titled

    untitled = build_figure_prompt(PageSource(page_number=4))
    assert "page 4" in untitled


def test_a_whole_answer_code_fence_is_unwrapped():
    """Models wrap Markdown in a fence despite being told not to; leaving it
    would put fence markers into the index and the notes prompt."""
    assert _strip_code_fence("```markdown\n# Title\n\nBody\n```") == "# Title\n\nBody"


def test_a_genuine_code_block_inside_a_page_keeps_its_fences():
    body = "Intro\n\n```python\nprint(1)\n```"
    assert _strip_code_fence(body) == body


def test_vision_image_base64_round_trips():
    import base64

    image = VisionImage(data=b"\x89PNGdata", media_type="image/png")
    assert base64.b64decode(image.to_base64()) == b"\x89PNGdata"
